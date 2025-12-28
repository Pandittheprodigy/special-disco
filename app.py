import os

# ==========================
# 1. CRITICAL: THREAD & SIGNAL FIXES
# ==========================
os.environ["CREWAI_DISABLE_TELEMETRY"] = "true"
os.environ["OTEL_SDK_DISABLED"] = "true"

import streamlit as st
import json
import sys
import io
import logging
from crewai import Agent, Task, Crew, Process, LLM
from crewai.tools import BaseTool
from langchain_community.tools import DuckDuckGoSearchRun
from pptx import Presentation

# Silence Streamlit's internal thread warnings
logging.getLogger('streamlit.runtime.scriptrunner_utils.script_run_context').setLevel(logging.ERROR)

# ==========================
# 2. UI CONFIGURATION
# ==========================
st.set_page_config(page_title="Ivy League Research Crew", layout="wide")

st.markdown("""
    <style>
    .log-container {
        background-color: #0d1117;
        color: #e6edf3;
        padding: 20px;
        border-radius: 10px;
        font-family: 'SF Mono', 'Fira Code', monospace;
        height: 450px;
        overflow-y: scroll;
        font-size: 0.85rem;
        border: 1px solid #30363d;
        line-height: 1.4;
    }
    .stButton>button { width: 100%; border-radius: 5px; height: 3em; background-color: #1f6feb; color: white; }
    </style>
""", unsafe_allow_html=True)

st.title("🏛️ Harvard Senior Research & Publication Directorate")
st.caption("High-Rigor Quantitative Analysis | 2025 Real-Time Empirical Data")

# --- Sidebar Configuration ---
def render_agent_config(agent_label, key_prefix):
    with st.sidebar.expander(f"🧬 {agent_label} Node", expanded=False):
        provider = st.selectbox("Provider", ["Google Gemini", "Groq", "OpenRouter"], key=f"{key_prefix}_p")
        api_key = st.text_input("API Key", type="password", key=f"{key_prefix}_k")
        
        # Logic Fix: Define the default BEFORE or INSIDE the conditional checks
        if provider == "Groq": 
            default_model = "llama-3.3-70b-versatile"
        elif provider == "OpenRouter": 
            default_model = "meta-llama/llama-3.3-70b-instruct:free"
        else
            # Default for Google Gemini or any other selection
            #default_model = "gemini-2.0-flash"
            model_name = st.text_input("Model Name", value=default_model, key=f"{key_prefix}_m")
        return provider, api_key, model_name
with st.sidebar:
    st.header("🔐 Access Credentials")
    st.info("Utilize high-parameter models (e.g., Llama 3.3 70B or Gemini Pro) for best academic results.")
    res_config = render_agent_config("Quantitative Researcher", "res")
    wri_config = render_agent_config("Lead Academic Editor", "wri")
    pre_config = render_agent_config("Presentation Principal", "pre")

# ==========================
# 3. ADVANCED ACADEMIC TOOLS
# ==========================

class WebSearchTool(BaseTool):
    name: str = "Web Search"
    description: str = "Search the internet for empirical data, 2025 statistics, and peer-reviewed summaries."
    def _run(self, query: str) -> str:
        return DuckDuckGoSearchRun().run(query)

class PowerPointCreatorTool(BaseTool):
    name: str = "Create PowerPoint"
    description: str = "Generates a professional academic .pptx file from a JSON schema."
    def _run(self, slides_json: str) -> str:
        try:
            clean_json = slides_json.replace("```json", "").replace("```", "").strip()
            start, end = clean_json.find('{'), clean_json.rfind('}') + 1
            data = json.loads(clean_json[start:end])
            
            prs = Presentation()
            for slide_info in data.get('slides', []):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                # Title formatting
                title_shape = slide.shapes.title
                title_shape.text = slide_info.get('title', 'Academic Analysis')
                
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.word_wrap = True
                    for point in slide_info.get('content', []):
                        p = tf.add_paragraph()
                        p.text = str(point)
                        p.level = 0
            
            prs.save("final_publication_deck.pptx")
            return "SUCCESS: Academic presentation generated."
        except Exception as e:
            return f"Error: {str(e)}"

# ==========================
# 4. ORCHESTRATION ENGINE
# ==========================

def create_crew_llm(config_tuple):
    provider, api_key, model_name = config_tuple
    if not api_key: return None
    prefix_map = {"Google Gemini": "gemini/", "Groq": "groq/", "OpenRouter": "openrouter/"}
    full_model = f"{prefix_map.get(provider, '')}{model_name}"
    # Setting temperature to 0.2 for strict factual adherence
    return LLM(model=full_model, api_key=api_key, temperature=0.2)

topic = st.text_area("Research Objective", placeholder="e.g., A quantitative analysis of global semiconductor supply chain resilience in Q4 2025")

if st.button("🚀 Commence Academic Synthesis"):
    llm_res = create_crew_llm(res_config)
    llm_wri = create_crew_llm(wri_config)
    llm_pre = create_crew_llm(pre_config)

    if not (llm_res and llm_wri and llm_pre):
        st.error("Credential validation failed. Ensure all API keys are active.")
    else:
        st.subheader("🕵️ Researcher Terminal (Live Output)")
        log_window = st.empty()
        
        class StreamToStreamlit:
            def __init__(self, display):
                self.display, self.buffer = display, ""
            def write(self, data):
                self.buffer += data
                self.display.markdown(f'<div class="log-container">{self.buffer.replace("\n", "<br>")}</div>', unsafe_allow_html=True)
            def flush(self): pass

        sys.stdout = StreamToStreamlit(log_window)

        try:
            # --- HIGH-RIGOR AGENTS ---
            researcher = Agent(
                role='Director of Quantitative Research',
                goal=f'Conduct an exhaustive empirical investigation into {topic}, extracting concrete facts, 2025 figures, and statistical data points.',
                backstory="""You are a tenured Principal Investigator with a PhD in Econometrics. 
                You have a zero-tolerance policy for speculation. Your work relies solely on hard data, 
                current 2025 market figures, and validated research. You cite all sources meticulously.""",
                tools=[WebSearchTool()],
                llm=llm_res, verbose=True
            )
            
            writer = Agent(
                role='Senior Editor - Harvard Business Review',
                goal='Synthesize raw empirical data into a rigorous, peer-review quality academic paper.',
                backstory="""You are an expert in scientific and business communication. You transform data 
                into compelling logical arguments. You ensure the paper includes a dedicated 'Facts & Figures' 
                executive summary and clear quantitative analysis throughout.""",
                llm=llm_wri, verbose=True
            )
            
            presenter = Agent(
                role='Principal Visual Communication Consultant',
                goal='Translate the complex academic paper into a data-driven PowerPoint presentation.',
                backstory="""You specialize in information design for CEO-level and Academic Board audiences. 
                You ensure every slide features a specific metric, fact, or figure. You avoid generic bullet 
                points in favor of high-density factual communication.""",
                tools=[PowerPointCreatorTool()],
                llm=llm_pre, verbose=True
            )

            # --- STRUCTURED TASKS ---
            t1 = Task(
                description=f"""Conduct a deep-dive research into {topic}. 
                1. Identify at least 5 key statistical facts or market figures from 2025. 
                2. Provide a detailed analysis of the current landscape. 
                3. Cite every major claim.""",
                expected_output="A data-dense research dossier with clearly labeled 'Key Metrics' and 'Empirical Findings'.",
                agent=researcher
            )
            
            t2 = Task(
                description="""Using the research dossier, draft a 1,500-word academic paper. 
                Structure: Abstract, Quantitative Landscape, Empirical Analysis, Strategic Recommendations, and Bibliography. 
                Ensure data points from the researcher are highlighted in bold.""",
                expected_output="A professional academic paper in Markdown format.",
                agent=writer,
                context=[t1]
            )
            
            t3 = Task(
                description="""Synthesize the paper into a 7-slide deck JSON. 
                - Slide 1: Title & Abstract
                - Slide 2: Methodology & Key Metrics
                - Slides 3-6: Quantitative Analysis (Include specific figures)
                - Slide 7: Conclusion.
                Pass this JSON to the Create PowerPoint tool.""",
                expected_output="Confirmation of .pptx file generation containing quantitative data.",
                agent=presenter,
                context=[t2]
            )

            crew = Crew(
                agents=[researcher, writer, presenter], 
                tasks=[t1, t2, t3], 
                process=Process.sequential,
                share_crew=False
            )
            
            crew.kickoff()
            sys.stdout = sys.__stdout__

            st.success("✅ Publication Process Concluded Successfully.")
            
            # --- TABBED RESULTS ---
            tab1, tab2 = st.tabs(["📄 Peer-Reviewed Paper", "📊 Data Export"])
            with tab1:
                st.markdown(t2.output.raw)
            with tab2:
                if os.path.exists("final_publication_deck.pptx"):
                    with open("final_publication_deck.pptx", "rb") as f:
                        st.download_button("📥 Download Quantitative Presentation (.pptx)", f, "harvard_analysis.pptx")

        except Exception as e:
            sys.stdout = sys.__stdout__
            st.error(f"System Error in Academic Pipeline: {str(e)}")
